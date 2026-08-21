import os
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ----------------------------------------------------
# 1. GENERATE THE 3 STEP-BY-STEP DIAGRAMS FOR SLIDES 4, 5, 6
# ----------------------------------------------------
def generate_flow_diagram(filename, step_num, title, direction, packet_role, flags_text, seq_text, client_state, server_state, color_theme):
    fig, ax = plt.subplots(figsize=(6.4, 5.4), dpi=300)
    ax.set_facecolor('#F8FAFC')
    fig.patch.set_facecolor('#F8FAFC')
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 7.2)
    ax.axis('off')

    # Step Title inside plot
    ax.text(5, 6.7, f'Step {step_num}: {title}', fontsize=12, fontweight='bold', ha='center', color='#0F172A', family='sans-serif')
    
    # 3 Main Nodes (Client, Router, Server)
    # Client Node (Role 1)
    ax.text(1.8, 5.2, 'CLIENT\n(Role 1: Host A)', fontsize=9.5, fontweight='bold', ha='center', va='center', color='#FFFFFF', family='sans-serif',
            bbox=dict(boxstyle='round,pad=0.5', facecolor='#2563EB', edgecolor='none'))
    ax.text(1.8, 4.1, f'State: {client_state}', fontsize=8, fontweight='bold', ha='center', color='#1E293B', family='sans-serif',
            bbox=dict(boxstyle='round,pad=0.3', facecolor='#EFF6FF', edgecolor='#BFDBFE', lw=1))

    # Router Node (Role 3)
    ax.text(5.0, 5.2, 'ROUTER\n(Role 3: Switch)', fontsize=9.5, fontweight='bold', ha='center', va='center', color='#FFFFFF', family='sans-serif',
            bbox=dict(boxstyle='round,pad=0.5', facecolor='#475569', edgecolor='none'))
    ax.text(5.0, 4.1, 'Packet Forwarding', fontsize=8, fontweight='bold', ha='center', color='#64748B', family='sans-serif')

    # Server Node (Role 4)
    ax.text(8.2, 5.2, 'SERVER\n(Role 4: Host B)', fontsize=9.5, fontweight='bold', ha='center', va='center', color='#FFFFFF', family='sans-serif',
            bbox=dict(boxstyle='round,pad=0.5', facecolor='#059669', edgecolor='none'))
    ax.text(8.2, 4.1, f'State: {server_state}', fontsize=8, fontweight='bold', ha='center', color='#1E293B', family='sans-serif',
            bbox=dict(boxstyle='round,pad=0.3', facecolor='#ECFDF5', edgecolor='#A7F3D0', lw=1))

    # Flow Arrows & Packet Motion
    if direction == 'LtoR':  # Client -> Router -> Server
        # Arrow 1: Client -> Router
        ax.annotate('', xy=(4.2, 2.7), xytext=(2.6, 2.7),
                    arrowprops=dict(arrowstyle='->', color=color_theme, lw=3, mutation_scale=16))
        # Arrow 2: Router -> Server
        ax.annotate('', xy=(7.4, 2.7), xytext=(5.8, 2.7),
                    arrowprops=dict(arrowstyle='->', color=color_theme, lw=3, mutation_scale=16))
        
        # Packet Role Badges along the path
        ax.text(3.4, 3.1, f'{packet_role}', fontsize=8.5, fontweight='bold', ha='center', va='center', color=color_theme, family='sans-serif')
        ax.text(6.6, 3.1, 'Forwarding to Server', fontsize=8.5, fontweight='bold', ha='center', va='center', color='#64748B', family='sans-serif')
    else:  # Server -> Router -> Client
        # Arrow 1: Server -> Router
        ax.annotate('', xy=(5.8, 2.7), xytext=(7.4, 2.7),
                    arrowprops=dict(arrowstyle='->', color=color_theme, lw=3, mutation_scale=16))
        # Arrow 2: Router -> Client
        ax.annotate('', xy=(2.6, 2.7), xytext=(4.2, 2.7),
                    arrowprops=dict(arrowstyle='->', color=color_theme, lw=3, mutation_scale=16))
        
        # Packet Role Badges along the path
        ax.text(6.6, 3.1, f'{packet_role}', fontsize=8.5, fontweight='bold', ha='center', va='center', color=color_theme, family='sans-serif')
        ax.text(3.4, 3.1, 'Routing to Client', fontsize=8.5, fontweight='bold', ha='center', va='center', color='#64748B', family='sans-serif')

    # Packet Parameters Box
    ax.text(5.0, 1.2, f'Flags: [{flags_text}]\n{seq_text}', fontsize=9, fontweight='bold', ha='center', va='center', color='#0F172A', family='sans-serif',
            bbox=dict(boxstyle='round,pad=0.5', facecolor='#FFFFFF', edgecolor='#CBD5E1', lw=1.5))

    plt.tight_layout()
    plt.savefig(filename, dpi=300, bbox_inches='tight')
    plt.close()
    print(f"Generated {filename}")

def build_all_diagrams():
    generate_flow_diagram('diagram_step1_syn.png', 1, 'SYN Packet Exchange', 'LtoR', 'SYN Packet (Role 2)', 'SYN = 1, ACK = 0', 'Seq = x (ISN_c) | Ack = 0', 'CLOSED -> SYN-SENT', 'LISTEN', '#2563EB')
    generate_flow_diagram('diagram_step2_synack.png', 2, 'SYN-ACK Packet Exchange', 'RtoL', 'SYN-ACK Packet (Role 5)', 'SYN = 1, ACK = 1', 'Seq = y (ISN_s) | Ack = x + 1', 'SYN-SENT -> ESTABLISHED', 'LISTEN -> SYN-RCVD', '#D97706')
    generate_flow_diagram('diagram_step3_ack.png', 3, 'ACK Packet Exchange', 'LtoR', 'ACK Packet (Role 6)', 'SYN = 0, ACK = 1', 'Seq = x + 1 | Ack = y + 1', 'ESTABLISHED', 'SYN-RCVD -> ESTABLISHED', '#059669')

# ----------------------------------------------------
# 2. HELPER FUNCTIONS FOR PPTX BUILDING
# ----------------------------------------------------
def add_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, tag_text, title_text):
    # Header tag badge
    tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
    tf = tag_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = tag_text.upper()
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = RGBColor(37, 99, 235)  # Accent Blue
    p.font.name = 'Trebuchet MS'

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.6))
    tf2 = title_box.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
    p2 = tf2.paragraphs[0]
    p2.text = title_text
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(15, 23, 42)  # Dark Navy
    p2.font.name = 'Trebuchet MS'

def add_card(slide, left, top, width, height, bg_color=RGBColor(255, 255, 255), border_color=RGBColor(226, 232, 240)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

# ----------------------------------------------------
# 3. PPTX DECK BUILDER
# ----------------------------------------------------
def build_presentation():
    build_all_diagrams()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Theme Colors
    BG_COLOR = RGBColor(248, 250, 252)       # Light Slate #F8FAFC
    NAVY_DARK = RGBColor(15, 23, 42)         # Deep Navy #0F172A
    CARD_BG = RGBColor(255, 255, 255)         # White
    CARD_BORDER = RGBColor(226, 232, 240)     # Light Border #E2E8F0
    ACCENT_BLUE = RGBColor(37, 99, 235)      # Royal Blue #2563EB
    ACCENT_GREEN = RGBColor(5, 150, 105)     # Emerald #059669
    ACCENT_AMBER = RGBColor(217, 119, 6)     # Amber #D97706
    TEXT_MUTED = RGBColor(71, 85, 105)       # Slate #475569

    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Role-Play Theme)
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    add_slide_background(slide1, NAVY_DARK)

    # Decorative top bar
    top_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = ACCENT_BLUE
    top_bar.line.fill.background()

    tb1 = slide1.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(4.5))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "COMPUTER NETWORKS | CLASSROOM ROLE-PLAY ACTIVITY"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = 'Trebuchet MS'
    p.space_after = Pt(14)

    p = tf1.add_paragraph()
    p.text = "TCP 3-Way Handshake"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = 'Trebuchet MS'
    p.space_after = Pt(10)

    p = tf1.add_paragraph()
    p.text = "Interactive Physical Role-Play: Client, Router, Server & Packet Transmission"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(203, 213, 225)
    p.font.name = 'Calibri'
    p.space_after = Pt(36)

    # 3 Role-Play Feature Cards
    pillars = [
        ("👥 6 Team Performers", "Client, Router, Server & 3 Packets"),
        ("🔀 Physical Packet Flow", "Client <-> Router <-> Server Traversal"),
        ("⚡ 3-Step Handshake", "SYN -> SYN-ACK -> ACK Synchronization")
    ]

    card_w = Inches(3.6)
    card_h = Inches(1.3)
    start_x = Inches(1.0)
    gap = Inches(0.26)

    for i, (title_p, desc_p) in enumerate(pillars):
        cx = start_x + i * (card_w + gap)
        add_card(slide1, cx, Inches(4.6), card_w, card_h, bg_color=RGBColor(30, 41, 59), border_color=RGBColor(51, 65, 85))
        
        tbox = slide1.shapes.add_textbox(cx + Inches(0.15), Inches(4.7), card_w - Inches(0.3), card_h - Inches(0.2))
        tf_c = tbox.text_frame
        tf_c.word_wrap = True
        
        p = tf_c.paragraphs[0]
        p.text = title_p
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = 'Trebuchet MS'
        p.space_after = Pt(4)
        
        p2 = tf_c.add_paragraph()
        p2.text = desc_p
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = RGBColor(148, 163, 184)
        p2.font.name = 'Calibri'

    tbox_ft = slide1.shapes.add_textbox(Inches(1.0), Inches(6.5), Inches(11.333), Inches(0.5))
    tf_ft = tbox_ft.text_frame
    p_ft = tf_ft.paragraphs[0]
    p_ft.text = "Classroom Activity Deck | Explainer Narrator + 6 Role-Play Performers | Computer Networks Course"
    p_ft.font.size = Pt(11)
    p_ft.font.color.rgb = RGBColor(100, 116, 139)
    p_ft.font.name = 'Calibri'

    slide1.notes_slide.notes_text_frame.text = (
        "Welcome class! Today, we are demonstrating the TCP 3-Way Handshake through an interactive 6-member network role-play activity. "
        "One narrator will explain the technical concepts while six team members physically enact the hosts, network router, and packet control flags moving across the classroom."
    )

    # ----------------------------------------------------
    # SLIDE 2: Overview & Role-Play Setup
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_slide_background(slide2, BG_COLOR)
    add_header(slide2, "01 | OVERVIEW", "Introduction & Role-Play Activity Concept")

    card_w = Inches(3.7)
    card_h = Inches(4.5)
    start_x = Inches(0.8)
    gap = Inches(0.3)

    cards_data_s2 = [
        ("Why a Physical Role-Play?", "Visualizing Protocol Mechanics", [
            ("Hop-by-Hop Visibility", "Shows packets traversing intermediate routers rather than instant abstract delivery."),
            ("Host State Tracking", "Demonstrates live state sign changes (CLOSED -> LISTEN -> ESTABLISHED)."),
            ("Interactive Learning", "Connects theoretical socket code to physical packet movement.")
        ], ACCENT_BLUE),
        ("Key Protocol Parameters", "Control Flags & Sequence Math", [
            ("SYN (Synchronize)", "Requests starting sequence synchronization from receiver."),
            ("ACK (Acknowledge)", "Confirms receipt of sequence bytes (Ack = Seq + 1)."),
            ("ISN (Initial Seq No)", "Random 32-bit sequence offset generated by Client (x) and Server (y).")
        ], ACCENT_GREEN),
        ("Role-Play Execution Rules", "Classroom Floor Instructions", [
            ("Packet Performers", "Physically carry flag cards and sequence numbers across the classroom floor."),
            ("Router Performer", "Inspects destination IP/Port and routes packets between hosts."),
            ("Host Performers", "Update state signs immediately upon receiving valid packets.")
        ], ACCENT_AMBER)
    ]

    for i, (ctitle, csub, bullets, header_color) in enumerate(cards_data_s2):
        cx = start_x + i * (card_w + gap)
        add_card(slide2, cx, Inches(1.5), card_w, card_h)

        card_hdr = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, Inches(1.5), card_w, Inches(0.75))
        card_hdr.fill.solid()
        card_hdr.fill.fore_color.rgb = header_color
        card_hdr.line.fill.background()

        tf_h = card_hdr.text_frame
        tf_h.word_wrap = True
        tf_h.margin_left = tf_h.margin_right = Inches(0.2)
        tf_h.margin_top = Inches(0.1)
        p = tf_h.paragraphs[0]
        p.text = ctitle
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = 'Trebuchet MS'

        p_sub = tf_h.add_paragraph()
        p_sub.text = csub
        p_sub.font.size = Pt(9.5)
        p_sub.font.color.rgb = RGBColor(241, 245, 249)
        p_sub.font.name = 'Calibri'

        c_tb = slide2.shapes.add_textbox(cx + Inches(0.2), Inches(2.35), card_w - Inches(0.4), card_h - Inches(0.95))
        tf_body = c_tb.text_frame
        tf_body.word_wrap = True

        for b_title, b_desc in bullets:
            p_b = tf_body.add_paragraph() if tf_body.paragraphs[0].text else tf_body.paragraphs[0]
            p_b.text = f"• {b_title}: "
            p_b.font.bold = True
            p_b.font.size = Pt(11)
            p_b.font.color.rgb = NAVY_DARK
            p_b.font.name = 'Calibri'

            run = p_b.add_run()
            run.text = b_desc
            run.font.bold = False
            run.font.size = Pt(10.5)
            run.font.color.rgb = TEXT_MUTED
            p_b.space_after = Pt(12)

    # Bottom Banner
    banner = add_card(slide2, Inches(0.8), Inches(6.2), Inches(11.733), Inches(0.8), bg_color=RGBColor(239, 246, 255), border_color=RGBColor(191, 219, 254))
    b_tb = slide2.shapes.add_textbox(Inches(1.0), Inches(6.25), Inches(11.333), Inches(0.7))
    tf_bn = b_tb.text_frame
    tf_bn.word_wrap = True
    p_bn = tf_bn.paragraphs[0]
    p_bn.text = "💡 Fundamental TCP Rule:"
    p_bn.font.bold = True
    p_bn.font.size = Pt(11.5)
    p_bn.font.color.rgb = ACCENT_BLUE
    p_bn.font.name = 'Trebuchet MS'
    
    run_bn = p_bn.add_run()
    run_bn.text = " Application data (e.g. HTTP requests) cannot flow until the 3-Way Handshake successfully transitions both Client and Server to the ESTABLISHED state."
    run_bn.font.bold = False
    run_bn.font.size = Pt(11)
    run_bn.font.color.rgb = NAVY_DARK

    slide2.notes_slide.notes_text_frame.text = (
        "Before we begin the physical demonstration, let's review why we perform a role-play. "
        "TCP is a stateful protocol operating hop-by-hop. In our classroom activity, packet performers will physically carry control flags across the network router. "
        "Both host performers will update their state signs as packets arrive."
    )

    # ----------------------------------------------------
    # SLIDE 3: Our Network Role-Play (The 6 Roles)
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_slide_background(slide3, BG_COLOR)
    add_header(slide3, "02 | CAST & ROLES", "Our Network Role-Play: 6 Team Roles")

    # 2 Rows x 3 Columns Grid for 6 Roles
    grid_w = Inches(3.7)
    grid_h = Inches(2.4)
    positions_s3 = [
        (Inches(0.8), Inches(1.5)),
        (Inches(4.816), Inches(1.5)),
        (Inches(8.833), Inches(1.5)),
        (Inches(0.8), Inches(4.15)),
        (Inches(4.816), Inches(4.15)),
        (Inches(8.833), Inches(4.15))
    ]

    roles_data = [
        ("👤 Role 1: Client / Sender", "Host A (Initiating Host)", [
            ("Function", "Initiates connection via socket connect()."),
            ("State Signs", "CLOSED -> SYN-SENT -> ESTABLISHED."),
            ("Sequence", "Generates Initial Seq Number (ISN_c = x).")
        ], ACCENT_BLUE),
        ("✉️ Role 2: SYN Packet", "Step 1 Packet Performer", [
            ("Function", "Carries initial connection request."),
            ("Flags & Parameters", "Flags: [SYN=1, ACK=0], Seq = x, Ack = 0."),
            ("Path", "Walks from Client -> Router -> Server.")
        ], ACCENT_BLUE),
        ("🔀 Role 3: Router / Switch", "Intermediate Network Hop", [
            ("Function", "Receives packets on ingress port & forwards to destination."),
            ("Inspection", "Verifies packet headers & routes across network."),
            ("Realism", "Demonstrates hop-by-hop network traversal.")
        ], RGBColor(71, 85, 105)),
        ("🖥️ Role 4: Server / Receiver", "Host B (Listening Host)", [
            ("Function", "Listens on port 443 & processes connection requests."),
            ("State Signs", "LISTEN -> SYN-RCVD -> ESTABLISHED."),
            ("Sequence", "Generates Server ISN (ISN_s = y).")
        ], ACCENT_GREEN),
        ("📩 Role 5: SYN-ACK Packet", "Step 2 Packet Performer", [
            ("Function", "Carries Server's response & sequence proposal."),
            ("Flags & Parameters", "Flags: [SYN=1, ACK=1], Seq = y, Ack = x+1."),
            ("Path", "Walks from Server -> Router -> Client.")
        ], ACCENT_AMBER),
        ("✅ Role 6: ACK Packet", "Step 3 Packet Performer", [
            ("Function", "Carries final confirmation to establish session."),
            ("Flags & Parameters", "Flags: [SYN=0, ACK=1], Seq = x+1, Ack = y+1."),
            ("Path", "Walks from Client -> Router -> Server.")
        ], ACCENT_GREEN)
    ]

    for i, (rtitle, rsub, rbullets, raccent) in enumerate(roles_data):
        rx, ry = positions_s3[i]
        add_card(slide3, rx, ry, grid_w, grid_h)

        stripe = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, rx, ry, Inches(0.12), grid_h)
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = raccent
        stripe.line.fill.background()

        tb_r = slide3.shapes.add_textbox(rx + Inches(0.22), ry + Inches(0.12), grid_w - Inches(0.3), grid_h - Inches(0.24))
        tf_r = tb_r.text_frame
        tf_r.word_wrap = True

        p_h = tf_r.paragraphs[0]
        p_h.text = rtitle
        p_h.font.size = Pt(12.5)
        p_h.font.bold = True
        p_h.font.color.rgb = NAVY_DARK
        p_h.font.name = 'Trebuchet MS'

        p_s = tf_r.add_paragraph()
        p_s.text = rsub.upper()
        p_s.font.size = Pt(8.5)
        p_s.font.bold = True
        p_s.font.color.rgb = raccent
        p_s.font.name = 'Trebuchet MS'
        p_s.space_after = Pt(4)

        for b_h, b_b in rbullets:
            p_b = tf_r.add_paragraph()
            p_b.text = f"• {b_h}: "
            p_b.font.bold = True
            p_b.font.size = Pt(9.5)
            p_b.font.color.rgb = NAVY_DARK

            r_b = p_b.add_run()
            r_b.text = b_b
            r_b.font.bold = False
            r_b.font.size = Pt(9)
            r_b.font.color.rgb = TEXT_MUTED
            p_b.space_after = Pt(2)

    # Footer note
    add_card(slide3, Inches(0.8), Inches(6.7), Inches(11.733), Inches(0.45), bg_color=RGBColor(241, 245, 249), border_color=RGBColor(203, 213, 225))
    tb_fn = slide3.shapes.add_textbox(Inches(1.0), Inches(6.72), Inches(11.333), Inches(0.4))
    tf_fn = tb_fn.text_frame
    p_fn = tf_fn.paragraphs[0]
    p_fn.text = "🎭 Performers Assigned — Ready to Begin 3-Step Physical Role-Play Demonstration!"
    p_fn.font.size = Pt(10.5)
    p_fn.font.bold = True
    p_fn.font.color.rgb = ACCENT_BLUE
    p_fn.font.name = 'Calibri'

    slide3.notes_slide.notes_text_frame.text = (
        "Let's introduce our 6 role-play performers! Role 1 is the Client Host. Role 2 is the SYN Packet. Role 3 is the Network Router. "
        "Role 4 is the Server Host. Role 5 is the SYN-ACK Packet. And Role 6 is the final ACK Packet. Each performer has a specific task during the 3-step exchange."
    )

    # ----------------------------------------------------
    # SLIDE 4: Step 1 — Client Connection Request (SYN Exchange)
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_slide_background(slide4, BG_COLOR)
    add_header(slide4, "03 | STEP 1: SYN", "Step 1: Client Connection Request (SYN Exchange)")

    # Left Side: Insert Diagram diagram_step1_syn.png (Width: 5.9)
    if os.path.exists("diagram_step1_syn.png"):
        slide4.shapes.add_picture("diagram_step1_syn.png", Inches(0.8), Inches(1.4), width=Inches(5.7))

    # Right Side: Role Performers Action Card (Width: 5.7)
    add_card(slide4, Inches(6.833), Inches(1.4), Inches(5.7), Inches(5.5))

    r_hdr4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.833), Inches(1.4), Inches(5.7), Inches(0.65))
    r_hdr4.fill.solid()
    r_hdr4.fill.fore_color.rgb = ACCENT_BLUE
    r_hdr4.line.fill.background()

    tf_rh4 = r_hdr4.text_frame
    tf_rh4.word_wrap = True
    p_rh4 = tf_rh4.paragraphs[0]
    p_rh4.text = "Physical Performer Actions (Step 1)"
    p_rh4.font.size = Pt(14)
    p_rh4.font.bold = True
    p_rh4.font.color.rgb = RGBColor(255, 255, 255)
    p_rh4.font.name = 'Trebuchet MS'

    tb_act4 = slide4.shapes.add_textbox(Inches(7.0), Inches(2.15), Inches(5.3), Inches(4.6))
    tf_act4 = tb_act4.text_frame
    tf_act4.word_wrap = True

    step1_actions = [
        ("Client (Role 1)", "Initiates Connection", "Calls connect(), generates ISN_c = x (e.g., 1000). Updates state sign: CLOSED -> SYN-SENT."),
        ("SYN Packet (Role 2)", "Dispatches Packet", "Holds card: [SYN=1, ACK=0], Seq = x, Ack = 0. Walks from Client to Router."),
        ("Router (Role 3)", "Forwards Packet", "Inspects SYN packet on ingress interface and forwards across classroom to Server."),
        ("Server (Role 4)", "Receives SYN", "In LISTEN state. Receives SYN, allocates half-open TCB buffer entry in SYN queue."),
        ("Sequence Rule", "Consumes 1 Seq Count", "SYN flag consumes 1 sequence number (x). Next expected byte will be x + 1.")
    ]

    for a_role, a_title, a_desc in step1_actions:
        p_a = tf_act4.add_paragraph() if tf_act4.paragraphs[0].text else tf_act4.paragraphs[0]
        p_a.text = f"• {a_role} [{a_title}]: "
        p_a.font.bold = True
        p_a.font.size = Pt(10.5)
        p_a.font.color.rgb = NAVY_DARK

        r_a = p_a.add_run()
        r_a.text = a_desc
        r_a.font.bold = False
        r_a.font.size = Pt(9.5)
        r_a.font.color.rgb = TEXT_MUTED
        p_a.space_after = Pt(6)

    slide4.notes_slide.notes_text_frame.text = (
        "Step 1 begins! Role 1, the Client, initiates an Active Open. Client hands the SYN Packet (Role 2) to the Router (Role 3). "
        "SYN Packet carries Seq=x. The Router forwards SYN Packet to the Server (Role 4). Client updates its state sign to SYN-SENT while Server stays in LISTEN."
    )

    # ----------------------------------------------------
    # SLIDE 5: Step 2 — Server Response & Sync (SYN-ACK Exchange)
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_slide_background(slide5, BG_COLOR)
    add_header(slide5, "04 | STEP 2: SYN-ACK", "Step 2: Server Response & Synchronization (SYN-ACK)")

    if os.path.exists("diagram_step2_synack.png"):
        slide5.shapes.add_picture("diagram_step2_synack.png", Inches(0.8), Inches(1.4), width=Inches(5.7))

    add_card(slide5, Inches(6.833), Inches(1.4), Inches(5.7), Inches(5.5))

    r_hdr5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.833), Inches(1.4), Inches(5.7), Inches(0.65))
    r_hdr5.fill.solid()
    r_hdr5.fill.fore_color.rgb = ACCENT_AMBER
    r_hdr5.line.fill.background()

    tf_rh5 = r_hdr5.text_frame
    tf_rh5.word_wrap = True
    p_rh5 = tf_rh5.paragraphs[0]
    p_rh5.text = "Physical Performer Actions (Step 2)"
    p_rh5.font.size = Pt(14)
    p_rh5.font.bold = True
    p_rh5.font.color.rgb = RGBColor(255, 255, 255)
    p_rh5.font.name = 'Trebuchet MS'

    tb_act5 = slide5.shapes.add_textbox(Inches(7.0), Inches(2.15), Inches(5.3), Inches(4.6))
    tf_act5 = tb_act5.text_frame
    tf_act5.word_wrap = True

    step2_actions = [
        ("Server (Role 4)", "Responds & Syncs", "Generates Server ISN_s = y (e.g., 5000) and sets Ack = x + 1. Updates state sign: LISTEN -> SYN-RCVD."),
        ("SYN-ACK Packet (Role 5)", "Dispatches Packet", "Holds card: [SYN=1, ACK=1], Seq = y, Ack = x + 1. Walks from Server to Router."),
        ("Router (Role 3)", "Routes Back", "Receives SYN-ACK from Server interface and routes back across classroom to Client."),
        ("Client (Role 1)", "Receives SYN-ACK", "Verifies Ack = x + 1. Immediately updates state sign: SYN-SENT -> ESTABLISHED!"),
        ("Milestone", "Client Connected", "Client host is now ready to transmit application payload after final ACK.")
    ]

    for a_role, a_title, a_desc in step2_actions:
        p_a = tf_act5.add_paragraph() if tf_act5.paragraphs[0].text else tf_act5.paragraphs[0]
        p_a.text = f"• {a_role} [{a_title}]: "
        p_a.font.bold = True
        p_a.font.size = Pt(10.5)
        p_a.font.color.rgb = NAVY_DARK

        r_a = p_a.add_run()
        r_a.text = a_desc
        r_a.font.bold = False
        r_a.font.size = Pt(9.5)
        r_a.font.color.rgb = TEXT_MUTED
        p_a.space_after = Pt(6)

    slide5.notes_slide.notes_text_frame.text = (
        "Now Step 2! Server (Role 4) acknowledges Client's SYN by setting Ack=x+1 and generates its own sequence number Seq=y. "
        "Server hands the SYN-ACK Packet (Role 5) to the Router (Role 3). Router delivers it to Client (Role 1). Upon receiving SYN-ACK, Client updates its state sign to ESTABLISHED!"
    )

    # ----------------------------------------------------
    # SLIDE 6: Step 3 — Client Final Confirmation (ACK Exchange)
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_slide_background(slide6, BG_COLOR)
    add_header(slide6, "05 | STEP 3: ACK", "Step 3: Client Confirms Connection (ACK Exchange)")

    if os.path.exists("diagram_step3_ack.png"):
        slide6.shapes.add_picture("diagram_step3_ack.png", Inches(0.8), Inches(1.4), width=Inches(5.7))

    add_card(slide6, Inches(6.833), Inches(1.4), Inches(5.7), Inches(5.5))

    r_hdr6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.833), Inches(1.4), Inches(5.7), Inches(0.65))
    r_hdr6.fill.solid()
    r_hdr6.fill.fore_color.rgb = ACCENT_GREEN
    r_hdr6.line.fill.background()

    tf_rh6 = r_hdr6.text_frame
    tf_rh6.word_wrap = True
    p_rh6 = tf_rh6.paragraphs[0]
    p_rh6.text = "Physical Performer Actions (Step 3)"
    p_rh6.font.size = Pt(14)
    p_rh6.font.bold = True
    p_rh6.font.color.rgb = RGBColor(255, 255, 255)
    p_rh6.font.name = 'Trebuchet MS'

    tb_act6 = slide6.shapes.add_textbox(Inches(7.0), Inches(2.15), Inches(5.3), Inches(4.6))
    tf_act6 = tb_act6.text_frame
    tf_act6.word_wrap = True

    step3_actions = [
        ("Client (Role 1)", "Sends Final ACK", "Sets Seq = x + 1 and Ack = y + 1. Clears SYN flag (SYN=0)."),
        ("ACK Packet (Role 6)", "Dispatches Packet", "Holds card: [SYN=0, ACK=1], Seq = x + 1, Ack = y + 1. Walks from Client to Router."),
        ("Router (Role 3)", "Forwards ACK", "Receives final ACK segment and delivers it to Server interface."),
        ("Server (Role 4)", "Completes Handshake", "Receives ACK, verifies Ack = y + 1. Updates state sign: SYN-RCVD -> ESTABLISHED!"),
        ("Outcome", "Bi-Directional Sync", "Both sockets are now ESTABLISHED. Application data transfer is unlocked!")
    ]

    for a_role, a_title, a_desc in step3_actions:
        p_a = tf_act6.add_paragraph() if tf_act6.paragraphs[0].text else tf_act6.paragraphs[0]
        p_a.text = f"• {a_role} [{a_title}]: "
        p_a.font.bold = True
        p_a.font.size = Pt(10.5)
        p_a.font.color.rgb = NAVY_DARK

        r_a = p_a.add_run()
        r_a.text = a_desc
        r_a.font.bold = False
        r_a.font.size = Pt(9.5)
        r_a.font.color.rgb = TEXT_MUTED
        p_a.space_after = Pt(6)

    slide6.notes_slide.notes_text_frame.text = (
        "Step 3 finalizes the handshake! Client (Role 1) hands the ACK Packet (Role 6) to Router (Role 3) with Seq=x+1 and Ack=y+1. "
        "Router forwards the ACK Packet to Server (Role 4). Server receives ACK and updates its state sign to ESTABLISHED! Both endpoints are now ESTABLISHED!"
    )

    # ----------------------------------------------------
    # SLIDE 7: Connection Established & Role-Play Summary
    # ----------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    add_slide_background(slide7, BG_COLOR)
    add_header(slide7, "06 | SUMMARY", "Connection Established & Role-Play Activity Summary")

    # 4 Cards Grid (2x2)
    grid_w = Inches(5.7)
    grid_h = Inches(2.2)
    positions_s7 = [
        (Inches(0.8), Inches(1.5)),
        (Inches(6.833), Inches(1.5)),
        (Inches(0.8), Inches(3.95)),
        (Inches(6.833), Inches(3.95))
    ]

    takeaways = [
        ("1. Physical Packet Traversal", "Client <-> Router <-> Server", [
            ("Router Mediation", "Packets physically traveled through intermediate router hop (Role 3)."),
            ("Flag Discipline", "SYN requested sync; SYN-ACK acknowledged & proposed; ACK confirmed."),
            ("Latency", "Complete 3-step physical exchange took 1 Round Trip Time (1 RTT).")
        ], ACCENT_BLUE),
        ("2. Sequence & Ack Number Math", "Byte Offset Agreement", [
            ("Client Offset", "Client ISN_c = x -> Server ACK = x + 1 -> Next Client Seq = x + 1."),
            ("Server Offset", "Server ISN_s = y -> Client ACK = y + 1 -> Next Server Seq = y + 1."),
            ("Sequence Cost", "SYN control flag consumed 1 sequence count without carrying data payload.")
        ], ACCENT_GREEN),
        ("3. Host State Machine Alignment", "Synchronized State Signs", [
            ("Client Path", "CLOSED -> SYN-SENT -> ESTABLISHED (Active Open)."),
            ("Server Path", "CLOSED -> LISTEN -> SYN-RCVD -> ESTABLISHED (Passive Open)."),
            ("Final State", "Both Client and Server hold 'ESTABLISHED' signs visible to classroom.")
        ], ACCENT_AMBER),
        ("4. Transition to Application Data", "Payload Transmission Unlocked", [
            ("Session Active", "Full-duplex reliable connection is now active between Host A and Host B."),
            ("Application Flow", "Client can now transmit HTTP GET request over open TCP socket."),
            ("Role-Play Wrap", "Applause for our 6 Role-Play Performers and Narrator!")
        ], RGBColor(124, 58, 237))
    ]

    for i, (t_title, t_sub, t_bullets, t_accent) in enumerate(takeaways):
        gx, gy = positions_s7[i]
        add_card(slide7, gx, gy, grid_w, grid_h)

        stripe = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, gx, gy, Inches(0.12), grid_h)
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = t_accent
        stripe.line.fill.background()

        tb_t = slide7.shapes.add_textbox(gx + Inches(0.3), gy + Inches(0.12), grid_w - Inches(0.4), grid_h - Inches(0.24))
        tf_t = tb_t.text_frame
        tf_t.word_wrap = True

        p_h = tf_t.paragraphs[0]
        p_h.text = t_title
        p_h.font.size = Pt(12.5)
        p_h.font.bold = True
        p_h.font.color.rgb = NAVY_DARK
        p_h.font.name = 'Trebuchet MS'

        p_s = tf_t.add_paragraph()
        p_s.text = t_sub.upper()
        p_s.font.size = Pt(8.5)
        p_s.font.bold = True
        p_s.font.color.rgb = t_accent
        p_s.font.name = 'Trebuchet MS'
        p_s.space_after = Pt(4)

        for b_h, b_b in t_bullets:
            p_b = tf_t.add_paragraph()
            p_b.text = f"• {b_h}: "
            p_b.font.bold = True
            p_b.font.size = Pt(9.5)
            p_b.font.color.rgb = NAVY_DARK

            r_b = p_b.add_run()
            r_b.text = b_b
            r_b.font.bold = False
            r_b.font.size = Pt(9)
            r_b.font.color.rgb = TEXT_MUTED
            p_b.space_after = Pt(2)

    # Presenter Wrap-Up Footer Box
    add_card(slide7, Inches(0.8), Inches(6.3), Inches(11.733), Inches(0.8), bg_color=NAVY_DARK, border_color=None)
    tb_wrap = slide7.shapes.add_textbox(Inches(1.0), Inches(6.35), Inches(11.333), Inches(0.7))
    tf_wrap = tb_wrap.text_frame
    tf_wrap.word_wrap = True
    
    p_w = tf_wrap.paragraphs[0]
    p_w.text = "🎓 Role-Play Activity Conclusion:"
    p_w.font.bold = True
    p_w.font.size = Pt(11.5)
    p_w.font.color.rgb = ACCENT_BLUE
    p_w.font.name = 'Trebuchet MS'

    r_w = p_w.add_run()
    r_w.text = " Our 6 role performers demonstrated how TCP establishes a reliable connection across a network router in 3 physical steps. By synchronizing sequence numbers x and y, both Client and Server reached the ESTABLISHED state. Great job team!"
    r_w.font.bold = False
    r_w.font.size = Pt(11)
    r_w.font.color.rgb = RGBColor(241, 245, 249)

    slide7.notes_slide.notes_text_frame.text = (
        "To summarize: Our 6 role performers demonstrated how TCP establishes a reliable connection across a network router in 3 physical steps. "
        "By synchronizing sequence numbers x and y, both Client and Server reached the ESTABLISHED state. Application data can now flow freely. Great job to our 6 performers!"
    )

    output_filename = "tcp_3way_handshake_presentation.pptx"
    prs.save(output_filename)
    print(f"Role-Play Presentation deck saved as '{output_filename}'")

if __name__ == "__main__":
    build_presentation()
